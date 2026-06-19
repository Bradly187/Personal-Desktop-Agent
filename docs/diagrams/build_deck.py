"""Build the Personal Desktop Agent study deck (.pptx).

Sections: Title · Agenda · Project Overview · State Machines · Data Model.
Diagrams are pre-rendered PNGs under docs/diagrams/{state,overview,db}/.
Re-run after re-rendering any PNG:  python docs/diagrams/build_deck.py
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = pathlib.Path(".")
STATE = ROOT / "docs/diagrams/state"
OV = ROOT / "docs/diagrams/overview"
DB = ROOT / "docs/diagrams/db"

# ── palette ────────────────────────────────────────────────────────────────
INK   = RGBColor(0x1A, 0x1A, 0x2E)
MUTE  = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG    = RGBColor(0xF7, 0xF8, 0xFA)
IPAD  = RGBColor(0x2D, 0x6C, 0xDF)   # blue
BACK  = RGBColor(0x7C, 0x3A, 0xED)   # purple
GREEN = RGBColor(0x0E, 0x9F, 0x6E)   # overview
AMBER = RGBColor(0xD9, 0x73, 0x0D)   # data model

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

HEADER_H = Inches(1.05)
FOOTER_H = Inches(0.45)
PAD = Inches(0.35)


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, l, t, w, h, fill=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.line.fill.background()
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.shadow.inherit = False
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return tb


def header(s, title, accent, badge):
    """White header bar + title + section badge. Returns body-top EMU."""
    bg(s, BG)
    rect(s, 0, 0, SW, HEADER_H, WHITE)
    rect(s, 0, HEADER_H, SW, Inches(0.025), accent)
    rect(s, Inches(0.45), Inches(0.30), Inches(0.10), Inches(0.46), accent)
    text(s, Inches(0.72), Inches(0.18), Inches(9.4), Inches(0.72),
         [[(title, 23, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    b = rect(s, Inches(10.5), Inches(0.33), Inches(2.4), Inches(0.40), accent)
    bf = b.text_frame
    bf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = bf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = badge
    br.font.size = Pt(11)
    br.font.bold = True
    br.font.color.rgb = WHITE
    br.font.name = "Segoe UI"
    return HEADER_H + Inches(0.20)


def footer(s, source):
    text(s, Inches(0.5), SH - FOOTER_H, Inches(9), Inches(0.4),
         [[(source, 11, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(9.5), SH - FOOTER_H, Inches(3.3), Inches(0.4),
         [[("Personal Desktop Agent", 10, MUTE, False)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def fit_image(s, img_path, top, avail_h):
    iw, ih = Image.open(img_path).size
    avail_w = SW - 2 * PAD
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    l = int(PAD + (avail_w - w) // 2)
    t = int(top + (avail_h - h) // 2)
    s.shapes.add_picture(str(img_path), Emu(l), Emu(t), Emu(w), Emu(h))


def img_slide(title, accent, badge, img_path, source, caption=None):
    s = slide()
    body_top = header(s, title, accent, badge)
    cap_h = Inches(0.5) if caption else Inches(0.0)
    if caption:
        text(s, Inches(0.6), body_top, Inches(12.1), Inches(0.45),
             [[(caption, 13, MUTE, False)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    top = body_top + cap_h
    avail_h = SH - top - FOOTER_H - Inches(0.1)
    fit_image(s, img_path, top, avail_h)
    footer(s, source)
    return s


def bullets_slide(title, accent, badge, paras, source, two_col=None):
    """paras: list of (level, text, bold). two_col: (left_paras, right_paras)."""
    s = slide()
    body_top = header(s, title, accent, badge)

    def render(paras, l, w):
        runs = []
        for (lvl, txt, bold) in paras:
            if lvl == 0:
                runs.append([("▸  ", 16, accent, True), (txt, 16, INK, bold)])
            elif lvl == -1:  # section label
                runs.append([(txt, 13, accent, True)])
            else:
                runs.append([("     – ", 13, MUTE, False), (txt, 13, RGBColor(0x3a,0x3f,0x4a), False)])
        text(s, l, body_top + Inches(0.15), w, SH - body_top - FOOTER_H - Inches(0.3),
             runs, space_after=9, line_spacing=1.05)

    if two_col:
        render(two_col[0], Inches(0.6), Inches(6.0))
        render(two_col[1], Inches(6.9), Inches(6.0))
    else:
        render(paras, Inches(0.7), Inches(12.0))
    footer(s, source)
    return s


def table_slide(title, accent, badge, headers, rows, source, col_w):
    s = slide()
    body_top = header(s, title, accent, badge)
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_w = sum(col_w)
    left = Emu(int((SW - Inches(tbl_w)) / 2))
    top = body_top + Inches(0.25)
    height = Inches(0.5 + 0.62 * len(rows))
    gtbl = s.shapes.add_table(nrows, ncols, left, top, Inches(tbl_w), height).table
    for j, cw in enumerate(col_w):
        gtbl.columns[j].width = Inches(cw)
    # header
    for j, htxt in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = accent
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xEF,0xF2,0xF6)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.bold = (j == 0); r.font.color.rgb = INK if j == 0 else RGBColor(0x3a,0x3f,0x4a)
            r.font.name = "Segoe UI"
    footer(s, source)
    return s


def divider(section_no, title, subtitle, accent):
    s = slide(); bg(s, INK)
    rect(s, 0, Inches(3.0), SW, Inches(0.04), accent)
    text(s, Inches(1), Inches(2.05), Inches(11.3), Inches(0.6),
         [[(f"SECTION {section_no}", 16, accent, True)]], align=PP_ALIGN.CENTER)
    text(s, Inches(1), Inches(2.55), Inches(11.3), Inches(1.0),
         [[(title, 38, WHITE, True)]], align=PP_ALIGN.CENTER)
    text(s, Inches(1), Inches(3.9), Inches(11.3), Inches(0.7),
         [[(subtitle, 15, MUTE, False)]], align=PP_ALIGN.CENTER)
    return s


# ════════════════════════════════════════════════════════════════════════
# TITLE
# ════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, INK)
rect(s, 0, Inches(3.05), SW, Inches(0.03), IPAD)
text(s, Inches(1), Inches(1.85), Inches(11.3), Inches(1.0),
     [[("Personal Desktop Agent", 40, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.8),
     [[("System Study Guide", 26, RGBColor(0x9C,0xB4,0xE8), False)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(4.35), Inches(11.3), Inches(0.6),
     [[("Multimodal accessibility desktop control  ·  iPad sensor hub + RTX 5090 inference", 14, MUTE, False)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
     [[("Architecture · 4-gate routing · 12 state machines · 42-table data model  ·  2026-06-13", 12, MUTE, False)]],
     align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# AGENDA
# ════════════════════════════════════════════════════════════════════════
s = slide()
body_top = header(s, "Agenda", IPAD, "Overview")
agenda = [
    ("1", "Project Overview", "What it is, the end-to-end pipeline, 4-gate routing, models, the agent kernel", GREEN),
    ("2", "State Machines", "12 machines — iPad sensors/UI (5) + PC agent kernel (7)", BACK),
    ("3", "Data Model", "Two-tier storage + 42-table agent.db grouped into 5 ER views", AMBER),
]
y = body_top + Inches(0.35)
for num, ttl, desc, ac in agenda:
    rect(s, Inches(0.9), y, Inches(0.7), Inches(0.7), ac)
    text(s, Inches(0.9), y, Inches(0.7), Inches(0.7), [[(num, 26, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.85), y - Inches(0.02), Inches(10.5), Inches(0.45),
         [[(ttl, 22, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.85), y + Inches(0.42), Inches(10.8), Inches(0.4),
         [[(desc, 14, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.5)
footer(s, "Study this deck top-to-bottom; each section's diagrams map 1:1 to source files")

# ════════════════════════════════════════════════════════════════════════
# SECTION 1 — PROJECT OVERVIEW
# ════════════════════════════════════════════════════════════════════════
divider("1", "Project Overview", "The mission, the pipeline, and how a command is routed", GREEN)

bullets_slide(
    "What Is This?", GREEN, "Overview",
    [
        (0, "Multimodal accessibility desktop control for a single user with rheumatoid arthritis", True),
        (1, "An iPad Pro is the sensor hub + primary touch surface; a Windows PC + RTX 5090 runs inference and executes desktop actions", False),
        (0, "Four control surfaces, one vocabulary", True),
        (1, "Voice · hand gesture · iPad tilt · direct touch  →  a 16-verb action vocabulary", False),
        (0, "Local-first, privacy-first", True),
        (1, "100% local inference (Ollama → vLLM); cloud (Anthropic API) only as an escalation fallback", False),
        (1, "Gate 0 forces sensitive input local; every tool call lands in an append-only, hash-chained audit log", False),
        (0, "Scale of the system today", True),
        (1, "16 action verbs · 6-level sensor fusion @ 60 Hz · 42-table operational DB · ~1,440 tests", False),
    ],
    "CLAUDE.md · specs/ipad-sensor-focus/requirements.md")

img_slide("End-to-End Architecture", GREEN, "Overview",
          OV / "A1-architecture.png",
          "core/ipad_bridge.py · core/fusion_engine.py · core/hybrid_coordinator.py · core/command_executor.py",
          caption="iPad sensors stream over WebSocket → fused at 60 Hz → routed → executed on the Windows desktop")

img_slide("A Real-World Interaction: \"Scroll Down\"", GREEN, "Overview",
          OV / "happy_path.png",
          "Example Interaction",
          caption="User rests hand on iPad and tilts forward while saying \"Scroll\" -> fused -> executed")

img_slide("How a Command Is Routed — 4 Gates", GREEN, "Overview",
          OV / "A2-gate-routing.png",
          "core/hybrid_coordinator.py",
          caption="Each command passes privacy → confidence → complexity → VRAM → latency gates; the deciding gate is logged")

table_slide(
    "Inference Tiers", GREEN, "Overview",
    ["Domain", "Model", "Notes"],
    [
        ["Command", "llama3.1:8b (4.6 GB)", "verb-first; ~190 ms warm p50"],
        ["Code / Plan", "qwen3-coder:30b", "thinking ON; DevAgent specialist"],
        ["Math", "deepseek-r1:8b", "chain-of-thought retained"],
        ["Vision", "qwen3-vl:30b", "UI target → pixel grounding"],
        ["General", "gemma3:27b / gemma4:12b", "resident slot, no eviction churn"],
        ["Cloud fallback", "claude-haiku-4-5 / opus-4-8", "command / dev; 10 s breaker"],
    ],
    "inference/local_inference.py · inference/model_router.py · RTX 5090 · 32 GB VRAM",
    col_w=[2.7, 4.2, 5.0])

bullets_slide(
    "Action Vocabulary — 16 Verbs", GREEN, "Overview",
    None,
    "core/command_executor.py · core/domain_classifier.py",
    two_col=(
        [
            (-1, "11 ACCESSIBILITY VERBS  (iPad sensor pipeline)", False),
            (0, "CLICK · MOUSEDOWN · MOUSEUP", False),
            (0, "SCROLL · TYPE · DICTATE", False),
            (0, "OPEN · CLOSE · HOTKEY", False),
            (0, "CLARIFY · SCREENSHOT", False),
            (1, "Verb-first format from llama3.1:8b", False),
        ],
        [
            (-1, "5 DEV-AGENT VERBS  (specialist models)", False),
            (0, "WRITE_FILE · RUN_TERMINAL", False),
            (0, "EXPLAIN · SEARCH_WEB", False),
            (0, "READ_SCREEN", False),
            (1, "Free-form; emitted via DevAgent", False),
            (1, "DomainClassifier picks the pipeline; CommandExecutor handles all 16", False),
        ],
    ))

bullets_slide(
    "The Agent Kernel — AIOS Primitives", GREEN, "Overview",
    [
        (0, "AccessibilityScheduler — priority queue, 5 tiers; accessibility/voice/gesture run concurrently, dev/background are semaphore-gated so they never starve the user", False),
        (0, "ResourceGovernor — pain-aware; on a flare it evicts heavy models, pauses the indexer + dev admission, and relaxes sensor thresholds (hysteresis 0.6 / 0.4)", False),
        (0, "Supervisor — one-for-one liveness watchdog; restarts dead loops under a bounded budget, latches FAILED and degrades gracefully (TTS warning)", False),
        (0, "CircuitBreaker — latches a down inference backend so it fast-fails instead of costing a full timeout per call", False),
        (0, "MemoryManager — a syscall façade over agent.db + SemanticMemory with schema-validated writes", False),
        (0, "Saga + goal_queue — durable, idempotent, rollback-safe autonomous execution", False),
        (-1, "→ Each of these is a state machine. The next section diagrams them.", False),
    ],
    "core/scheduler.py · resource_governor.py · supervisor.py · circuit_breaker.py · storage/memory_manager.py")

# ════════════════════════════════════════════════════════════════════════
# SECTION 2 — STATE MACHINES
# ════════════════════════════════════════════════════════════════════════
divider("2", "State Machines", "12 machines — how each subsystem behaves over time", BACK)

state_slides = [
    ("01-ipadapp-top-level-application-lifecycle", "App Top-Level Lifecycle", "iPad", "DesktopAgentApp · SensorManager"),
    ("02-websocketmanager-connection-state-machin", "WebSocket Connection", "iPad", "WebSocketManager.swift"),
    ("03-tiltsensor-navigation-state-machine", "Tilt Navigation", "iPad", "TiltSensor.swift (Core Motion)"),
    ("04-keywordlistener-recognition-state-machin", "Keyword Listener", "iPad", "KeywordListener.swift (Speech)"),
    ("05-appmode-ui-view-state-machine", "UI View Mode", "iPad", "SwiftUI views"),
    ("06-fusionengine-6-level-priority-tick", "FusionEngine — 6-Level Tick", "Backend", "core/fusion_engine.py"),
    ("07-gyrobiascalibrator-tilt-drift-compensati", "Gyro Bias Calibrator", "Backend", "calibration/gyro_bias_calibrator.py"),
    ("08-circuitbreaker-inference-backend-latch", "Circuit Breaker", "Backend", "core/circuit_breaker.py"),
    ("09-resourcegovernor-pain-aware-flare-mode", "Resource Governor (Flare)", "Backend", "core/resource_governor.py"),
    ("10-supervisor-per-subsystem-restart-policy", "Supervisor Restart Policy", "Backend", "core/supervisor.py"),
    ("11-goal-queue-row-lifecycle", "Goal-Queue Row Lifecycle", "Backend", "storage/db.py · goal_queue"),
    ("12-agent-run-lifecycle", "Agent-Run Lifecycle", "Backend", "storage/db.py · agent_runs"),
]
n_ipad = sum(1 for x in state_slides if x[2] == "iPad")
i_ipad = i_back = 0
for stem, ttl, layer, src in state_slides:
    accent = IPAD if layer == "iPad" else BACK
    if layer == "iPad":
        i_ipad += 1; badge = f"iPad · Swift   {i_ipad}/{n_ipad}"
    else:
        i_back += 1; badge = f"PC Backend   {i_back}/{len(state_slides)-n_ipad}"
    img_slide(ttl, accent, badge, STATE / f"{stem}.png", f"Source: {src}")

# ════════════════════════════════════════════════════════════════════════
# SECTION 3 — DATA MODEL
# ════════════════════════════════════════════════════════════════════════
divider("3", "Data Model", "Two-tier storage and the 42-table operational schema", AMBER)

img_slide("Two-Tier Storage", AMBER, "Data Model",
          DB / "D1-storage-overview.png",
          "storage/db.py · storage/audit_log.py · storage/semantic_memory.py",
          caption="SQLite on the hot path, DuckDB for analytics (zero-ETL attach), ChromaDB for RAG, append-only audit.db")

db_slides = [
    ("D2-core-star-schema", "Core Star Schema", "commands is the central fact table; sessions is the anchor",
     "storage/db.py — sessions · commands · inferences · sensor_events · session_summaries"),
    ("D3a-orchestration-queue", "Dev-Agent Orchestration (Queue)", "Durable goal queue → runs → steps",
     "storage/db.py — goal_queue · agent_runs · agent_steps"),
    ("D3b-orchestration-capabilities", "Dev-Agent Orchestration (Capabilities)", "Saga rollback, escalations, tools, events, skills",
     "storage/db.py — saga_compensations · dev_escalations · tool_calls · event_rules · skill_invocations"),
    ("D4-learning-adaptation", "Learning & Adaptation", "Few-shot examples/counterexamples, gesture calibration, threshold history",
     "storage/db.py — few_shot_* · gesture_* · settings_versions · adaptation_log"),
    ("D5-twin-voice", "Behavioral Twin & Voice", "Pain-day signals, session history, and per-condition voice calibration",
     "storage/db.py — twin_* · voice_* · sensor_rom · flare_profile"),
    ("D6-telemetry-events", "Telemetry, Events & Limits", "1 Hz sensor snapshots, the event bus, and rate-limit/observability config",
     "storage/db.py — sensor_telemetry · event_log · event_consumers · rate_limit_* · tool_*_config"),
]
for stem, ttl, cap, src in db_slides:
    img_slide(ttl, AMBER, "Data Model", DB / f"{stem}.png", src, caption=cap)

# closing
s = slide(); bg(s, INK)
rect(s, 0, Inches(3.0), SW, Inches(0.03), IPAD)
text(s, Inches(1), Inches(2.2), Inches(11.3), Inches(0.8),
     [[("Where Everything Lives", 30, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(3.4), Inches(10.3), Inches(2.5),
     [
        [("Diagrams (Mermaid source + PNG):  ", 14, RGBColor(0x9C,0xB4,0xE8), True), ("docs/diagrams/{state,overview,db}/", 14, WHITE, False)],
        [("State-machine spec:  ", 14, RGBColor(0x9C,0xB4,0xE8), True), ("specs/ipad-sensor-focus/diagrams/04-state-machines.md", 14, WHITE, False)],
        [("DB design rationale:  ", 14, RGBColor(0x9C,0xB4,0xE8), True), ("docs/architecture/database-design.md", 14, WHITE, False)],
        [("Project orientation:  ", 14, RGBColor(0x9C,0xB4,0xE8), True), ("CLAUDE.md  (status, key files, conventions)", 14, WHITE, False)],
        [("Rebuild this deck:  ", 14, RGBColor(0x9C,0xB4,0xE8), True), ("python docs/diagrams/build_deck.py", 14, WHITE, False)],
     ], space_after=12)

out = ROOT / "docs/diagrams/agent-study-deck.pptx"
prs.save(out)
print("saved", out, "· slides:", len(prs.slides._sldIdLst))
